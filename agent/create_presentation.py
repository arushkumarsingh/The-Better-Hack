from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os
import json
from typing import List, Dict
from agent.create_google_presentation import create_google_feature_presentation

def create_feature_presentation(keyframe_summaries: List[str], 
                              user_journey: str,
                              image_paths: List[str],
                              output_path: str = "output/presentation",
                              language: str = None,
                              website_context: str = None):
    """
    Creates a professional PowerPoint presentation highlighting the main features
    with proper text wrapping and formatting
    """
    prs = Presentation()
    
    # Define slide dimensions
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)  # 16:9 aspect ratio
    
    # Set theme colors
    theme_color_primary = RGBColor(0, 112, 192)     # Blue
    theme_color_secondary = RGBColor(255, 255, 255) # White
    theme_color_accent = RGBColor(240, 240, 240)    # Light gray
    
    # --- Personalization Agent: Analyze website context and user journey ---
    personalization = None
    if website_context:
        try:
            import openai
            personalization_prompt = f"""
You are an expert product marketer and user experience researcher. Given the following website context and user journey, summarize:
1. What the website offers (value proposition)
2. What the user is likely looking for (user needs)
3. How to tailor the messaging and feature highlights in a product demo deck to best resonate with this user and their context.
Return a JSON object with keys: value_proposition, user_needs, personalized_messaging, personalized_titles (array of 5 for features).

Website Context:
{website_context}

User Journey:
{user_journey}
"""
            response = openai.chat.completions.create(
                model="gpt-4o",
                messages=[{"role": "user", "content": personalization_prompt}],
                response_format={ "type": "json_object" }
            )
            personalization = json.loads(response.choices[0].message.content)
        except Exception as e:
            personalization = None
    else:
        personalization = None

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]

    # Use personalized messaging if available
    if personalization and personalization.get("personalized_messaging"):
        subtitle.text = personalization["personalized_messaging"][:200]
        if len(personalization["personalized_messaging"]) > 200:
            extra_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
            extra_box.text = personalization["personalized_messaging"][200:600]
    elif website_context:
        subtitle.text = website_context[:200]
        if len(website_context) > 200:
            extra_box = title_slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(2))
            extra_box.text = website_context[200:600]
    
    # Localize static text if language is specified and not English
    localized_title = "Application Features Overview"
    localized_subtitle = "Generated from User Journey Analysis"
    localized_summary = "Key Features Summary"
    if language and language.lower() != "english":
        import openai
        translation_prompt = f"Translate the following phrases into {language}. Return a JSON object with keys 'title', 'subtitle', 'summary'.\nPhrases: title='Application Features Overview', subtitle='Generated from User Journey Analysis', summary='Key Features Summary'"
        response = openai.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{"role": "user", "content": translation_prompt}],
            max_tokens=150
        )
        import json
        try:
            translations = json.loads(response.choices[0].message.content.strip())
            localized_title = translations.get('title', localized_title)
            localized_subtitle = translations.get('subtitle', localized_subtitle)
            localized_summary = translations.get('summary', localized_summary)
        except Exception:
            pass
    # Format title slide
    title.text = localized_title
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = theme_color_primary

    subtitle.text = localized_subtitle
    subtitle.text_frame.paragraphs[0].font.size = Pt(24)
    subtitle.text_frame.paragraphs[0].font.italic = True
    
    # Add background shape to title slide
    background = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(4), prs.slide_width, Inches(1.625)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = theme_color_primary
    background.line.fill.background()  # No outline
    
    # Extract main features from user journey
    features = _extract_main_features(user_journey, language=language)
    # If we have personalized titles, apply them to features
    if personalization and personalization.get("personalized_titles"):
        for i, f in enumerate(features):
            if i < len(personalization["personalized_titles"]):
                f["title"] = personalization["personalized_titles"][i]
    
    # Create feature slides
    for i, (feature, image_path) in enumerate(zip(features[:5], image_paths[:5])):
        feature_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add feature number for visual interest
        feature_num = feature_slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(1), Inches(1)
        )
        feature_num.fill.solid()
        feature_num.fill.fore_color.rgb = theme_color_primary
        feature_num.line.fill.background()  # No outline
        
        # Add number text to the oval
        num_textframe = feature_num.text_frame
        num_textframe.text = str(i+1)
        num_textframe.paragraphs[0].alignment = PP_ALIGN.CENTER
        num_textframe.paragraphs[0].font.size = Pt(24)
        num_textframe.paragraphs[0].font.bold = True
        num_textframe.paragraphs[0].font.color.rgb = theme_color_secondary
        num_textframe.vertical_anchor = MSO_ANCHOR.MIDDLE
        
        # Add title with colored background
        title_shape = feature_slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(1.7), Inches(0.5), Inches(7.8), Inches(1)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = theme_color_accent
        title_shape.line.fill.background()  # No outline
        
        # Add title text
        title_textbox = feature_slide.shapes.add_textbox(
            Inches(1.8), Inches(0.65), Inches(7.5), Inches(0.7)
        )
        title_textframe = title_textbox.text_frame
        title_para = title_textframe.add_paragraph()
        title_para.text = feature['title']
        title_para.font.size = Pt(32)
        title_para.font.bold = True
        title_para.font.color.rgb = theme_color_primary
        
        # Add image
        if os.path.exists(image_path):
            left = Inches(1)
            top = Inches(1.8)
            pic = feature_slide.shapes.add_picture(
                image_path, 
                left, 
                top, 
                width=Inches(4)
            )
        
        # Add description with proper text wrapping
        textbox = feature_slide.shapes.add_textbox(
            left=Inches(5.5), 
            top=Inches(1.8),
            width=Inches(3.5),
            height=Inches(3)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True
        
        # Add description paragraph with proper formatting
        p = text_frame.add_paragraph()
        p.text = feature['description']
        p.font.size = Pt(16)
        p.alignment = PP_ALIGN.LEFT
        text_frame.auto_size = True  # Auto-size text to fit in textbox
        
    # Create summary slide
    summary_slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title
    summary_title = summary_slide.shapes.add_textbox(
        Inches(1), Inches(0.5), Inches(8), Inches(1)
    )
    summary_title_frame = summary_title.text_frame
    summary_para = summary_title_frame.add_paragraph()
    summary_para.text = localized_summary
    summary_para.font.size = Pt(40)
    summary_para.font.bold = True
    summary_para.font.color.rgb = theme_color_primary
    
    # Add summary list
    summary_textbox = summary_slide.shapes.add_textbox(
        Inches(1), Inches(1.7), Inches(8), Inches(3)
    )
    summary_frame = summary_textbox.text_frame
    summary_frame.word_wrap = True
    
    # If personalized value prop/user needs, show them at the top of the summary slide
    if personalization:
        if personalization.get("value_proposition"):
            vp_para = summary_frame.add_paragraph()
            vp_para.text = f"Value Proposition: {personalization['value_proposition'][:400]}"
            vp_para.font.size = Pt(14)
            vp_para.level = 0
            vp_para.space_after = Pt(8)
        if personalization.get("user_needs"):
            un_para = summary_frame.add_paragraph()
            un_para.text = f"User Needs: {personalization['user_needs'][:400]}"
            un_para.font.size = Pt(14)
            un_para.level = 0
            un_para.space_after = Pt(8)
    elif website_context:
        context_para = summary_frame.add_paragraph()
        context_para.text = f"Context: {website_context[:400]}"
        context_para.font.size = Pt(14)
        context_para.level = 0
        context_para.space_after = Pt(8)
    
    # Add each feature as a bullet point
    for i, feature in enumerate(features[:5]):
        bullet_para = summary_frame.add_paragraph()
        bullet_para.text = f"{feature['title']}: {feature['description'].split('.')[0]}."
        bullet_para.font.size = Pt(18)
        bullet_para.level = 0  # First level bullet
        bullet_para.space_after = Pt(12)  # Space between bullets
    
    # Create output directory if it doesn't exist
    os.makedirs(output_path, exist_ok=True)
    
    # Save presentation
    output_file = f"{output_path}/feature_overview.pptx"
    prs.save(output_file)
    return output_file

def _extract_main_features(user_journey: str, language: str = None) -> List[Dict]:
    """
    Extracts main features from user journey text
    Returns list of dicts with 'title' and 'description'
    """
    # Use OpenAI to extract main features
    import openai
    
    prompt = """
    Given this user journey, identify the top 5 main features of the application.
    For each feature provide:
    1. A short title (2-4 words)
    2. A brief description (2-3 sentences)
    Format as JSON list of objects with 'title' and 'description' fields.
    
    User Journey:
    """
    if language and language.lower() != "english":
        prompt += f"\nOutput ONLY in {language}."
    
    response = openai.chat.completions.create(
        model="gpt-4o",
        messages=[
            {"role": "user", 
             "content": prompt + user_journey}
        ],
        response_format={ "type": "json_object" }
    )
    
    # Parse the JSON string into a Python dictionary
    features = json.loads(response.choices[0].message.content)
    return features['features']